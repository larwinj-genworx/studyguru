from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from src.schemas.study_material import ConceptContentPack

if TYPE_CHECKING:
    from pptx import Presentation

class PptxRenderer:
    def render(
        self,
        output_dir: Path,
        subject_name: str,
        grade_level: str,
        concept_packs: list[ConceptContentPack],
    ) -> Path:
        from pptx import Presentation

        presentation = Presentation()
        self._add_title_slide(presentation, subject_name, grade_level, concept_packs)

        for pack in concept_packs:
            self._add_concept_slides(presentation, pack)

        output_path = output_dir / "study_material.pptx"
        presentation.save(str(output_path))
        return output_path

    @staticmethod
    def _add_title_slide(
        presentation: Any,
        subject_name: str,
        grade_level: str,
        concept_packs: list[ConceptContentPack],
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = f"{subject_name} - Study Material"
        slide.placeholders[1].text = (
            f"Level: {grade_level}\n"
            f"Concepts covered: {', '.join(pack.concept_name for pack in concept_packs)}"
        )

    def _add_concept_slides(self, presentation: Any, pack: ConceptContentPack) -> None:
        # Dynamic 5-8 slides per concept based on content depth.
        complexity = len(pack.key_steps) + len(pack.examples) + len(pack.common_mistakes)
        extra_slides = 0
        if complexity > 10:
            extra_slides += 1
        if len(pack.mcqs) >= 8:
            extra_slides += 1
        extra_slides = min(extra_slides, 2)

        self._add_bullet_slide(
            presentation,
            title=f"{pack.concept_name}: What and Why",
            bullets=[pack.definition, pack.intuition],
        )
        self._add_bullet_slide(
            presentation,
            title=f"{pack.concept_name}: Key Steps",
            bullets=pack.key_steps[:6],
        )
        self._add_bullet_slide(
            presentation,
            title=f"{pack.concept_name}: Worked Examples",
            bullets=pack.examples[:5],
        )
        self._add_bullet_slide(
            presentation,
            title=f"{pack.concept_name}: Common Mistakes",
            bullets=pack.common_mistakes[:5],
        )
        self._add_bullet_slide(
            presentation,
            title=f"{pack.concept_name}: Quick Recap",
            bullets=pack.recap[:6],
        )

        if extra_slides >= 1:
            practice_items = [item["question"] for item in pack.mcqs[:4]]
            self._add_bullet_slide(
                presentation,
                title=f"{pack.concept_name}: Practice Check",
                bullets=practice_items,
            )
        if extra_slides >= 2:
            refs = [
                f'{ref.get("title", "Resource")} - {ref.get("url", "")}'
                for ref in pack.references[:4]
            ]
            if not refs:
                refs = ["Use textbook examples and class notes for additional practice."]
            self._add_bullet_slide(
                presentation,
                title=f"{pack.concept_name}: Learning Resources",
                bullets=refs,
            )

    @staticmethod
    def _add_bullet_slide(presentation: Any, title: str, bullets: list[str]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for idx, bullet in enumerate(bullets):
            bullet = bullet.strip()
            if not bullet:
                continue
            if idx == 0:
                text_frame.text = bullet
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet
